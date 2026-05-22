import base64
import sys
from pathlib import Path
from typing import Generator

sys.path.insert(0, str(Path(__file__).parent.parent.parent))


def extract_chunks(file_path: str) -> Generator[dict, None, None]:
    """파일에서 페이지/슬라이드별 텍스트 청크를 추출합니다."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pptx":
        yield from _extract_pptx(path)
    elif suffix == ".pdf":
        yield from _extract_pdf(path)
    elif suffix == ".docx":
        yield from _extract_docx(path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {suffix}")


def _describe_images_with_claude(images: list[dict]) -> str:
    """이미지 목록을 Claude Vision으로 분석해 설명 텍스트를 반환합니다."""
    import config
    import anthropic

    if not config.ANTHROPIC_API_KEY:
        return ""

    client = anthropic.Anthropic(api_key=config.ANTHROPIC_API_KEY)

    content = []
    for img in images:
        content.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": img["media_type"],
                "data": img["data"],
            },
        })
    content.append({
        "type": "text",
        "text": (
            "이 슬라이드의 이미지들을 분석해주세요. "
            "차트, 다이어그램, 표, 텍스트, 아키텍처 구조 등 핵심 내용을 한국어로 요약해주세요. "
            "검색에 활용될 설명이므로 주요 키워드와 개념을 포함해주세요."
        ),
    })

    response = client.messages.create(
        model=config.VISION_MODEL,
        max_tokens=1024,
        messages=[{"role": "user", "content": content}],
    )
    return response.content[0].text.strip()


def _extract_pptx(path: Path) -> Generator[dict, None, None]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        images = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

            # 이미지 추출
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    media_type = img.content_type or "image/png"
                    images.append({
                        "media_type": media_type,
                        "data": base64.standard_b64encode(img.blob).decode("utf-8"),
                    })
                except Exception:
                    pass

        # 슬라이드 노트
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[노트] {notes}")

        # 이미지가 있으면 Claude Vision으로 설명 생성
        if images:
            vision_desc = _describe_images_with_claude(images)
            if vision_desc:
                texts.append(f"[이미지 설명] {vision_desc}")

        if texts:
            yield {
                "page": slide_num,
                "text": "\n".join(texts),
                "page_label": f"슬라이드 {slide_num}",
            }


def _extract_pdf(path: Path) -> Generator[dict, None, None]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            yield {
                "page": page_num,
                "text": text,
                "page_label": f"{page_num}페이지",
            }


def _extract_docx(path: Path) -> Generator[dict, None, None]:
    from docx import Document

    doc = Document(str(path))
    chunk_texts = []
    chunk_num = 1
    char_count = 0
    CHARS_PER_CHUNK = 1000

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        chunk_texts.append(text)
        char_count += len(text)

        if char_count >= CHARS_PER_CHUNK:
            yield {
                "page": chunk_num,
                "text": "\n".join(chunk_texts),
                "page_label": f"섹션 {chunk_num}",
            }
            chunk_num += 1
            chunk_texts = []
            char_count = 0

    if chunk_texts:
        yield {
            "page": chunk_num,
            "text": "\n".join(chunk_texts),
            "page_label": f"섹션 {chunk_num}",
        }
